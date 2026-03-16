#!/usr/bin/env python3
"""Generate the Parent Setup Guide PowerPoint deck (light theme with images)."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ASSETS = os.path.join(os.path.dirname(__file__), "deck_assets")
OUTPUT = os.path.join(os.path.dirname(__file__), "Parent_Setup_Guide.pptx")

BG_WHITE = RGBColor(0xFA, 0xFA, 0xFC)
BG_SOFT = RGBColor(0xF0, 0xF4, 0xF8)
TEXT_DARK = RGBColor(0x1A, 0x20, 0x2C)
TEXT_MED = RGBColor(0x4A, 0x50, 0x68)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)
ACCENT_TEAL = RGBColor(0x0D, 0x9B, 0x8C)
ACCENT_ORANGE = RGBColor(0xE8, 0x6E, 0x17)
ACCENT_GREEN = RGBColor(0x16, 0x87, 0x3D)
LINK_BLUE = RGBColor(0x1D, 0x4E, 0xD8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color=BG_WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text(slide, left, top, w, h, txt, size=18, bold=False,
         color=TEXT_DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tf


def bullets(slide, left, top, w, h, items, size=20, color=TEXT_DARK,
            spacing=Pt(10)):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
    return tf


def img(slide, name, left, top, width=None, height=None):
    path = os.path.join(ASSETS, name)
    if not os.path.exists(path):
        return
    kw = {"left": left, "top": top}
    if width:
        kw["width"] = width
    if height:
        kw["height"] = height
    slide.shapes.add_picture(path, **kw)


def link_line(tf, label, url, size=20):
    p = tf.add_paragraph()
    p.space_after = Pt(14)
    r1 = p.add_run()
    r1.text = f"{label}:  "
    r1.font.size = Pt(size)
    r1.font.color.rgb = TEXT_DARK
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = url
    r2.font.size = Pt(size)
    r2.font.color.rgb = LINK_BLUE
    r2.font.name = "Calibri"
    r2.font.underline = True
    r2.hyperlink.address = url


def accent_bar(slide, left, top, width, color=ACCENT_BLUE):
    bar = slide.shapes.add_shape(1, left, top, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── SLIDE 1: COVER ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_WHITE)
    text(s, Inches(0.8), Inches(1.2), Inches(7), Inches(1.5),
         "AI Agents Course", size=52, bold=True, color=TEXT_DARK)
    accent_bar(s, Inches(0.8), Inches(2.8), Inches(3), ACCENT_BLUE)
    text(s, Inches(0.8), Inches(3.2), Inches(7), Inches(1),
         "Parent / Guardian Setup Guide", size=26, color=TEXT_MED)
    text(s, Inches(0.8), Inches(4.3), Inches(7), Inches(1),
         "Everything you need to get your student ready\n~20 minutes total",
         size=20, color=ACCENT_BLUE)
    img(s, "parent_deck_cover.png", Inches(7.8), Inches(0.8), width=Inches(5))

    # ── SLIDE 2: WHAT IS THIS COURSE? ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_WHITE)
    text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.8),
         "What is This Course?", size=36, bold=True)
    accent_bar(s, Inches(0.6), Inches(1.1), Inches(2.5), ACCENT_BLUE)
    bullets(s, Inches(0.6), Inches(1.5), Inches(7), Inches(5), [
        "An 8-session hands-on course for high school students",
        "Students learn what AI is, how it works, and build their own AI agent",
        "The agent lives on a cloud server and talks to them via Telegram",
        "It can search the web, browse websites, and remember things",
        "No coding experience required — students customize with plain English",
        "",
        "By the end, each student presents a working AI agent they designed",
    ], size=22)
    img(s, "parent_deck_3steps.png", Inches(8), Inches(1.5), width=Inches(4.8))

    # ── SLIDE 3: STUDENT MACHINE PREREQUISITES ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_SOFT)
    text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.8),
         "What Your Student Needs", size=36, bold=True)
    accent_bar(s, Inches(0.6), Inches(1.1), Inches(2.5), ACCENT_TEAL)

    bullets(s, Inches(0.6), Inches(1.5), Inches(6.8), Inches(5), [
        "A laptop or desktop computer",
        "      Windows 10/11, macOS, or Linux — any is fine",
        "      Chromebook will NOT work",
        "",
        "VS Code (free code editor — install before class)",
        "      Download at code.visualstudio.com",
        "      Students will use this to edit their agent's files",
        "",
        "Telegram on their phone",
        "      iOS or Android — free",
        "",
        "Stable internet connection",
        "      Needed to connect to the cloud server",
    ], size=21, spacing=Pt(5))

    img(s, "parent_deck_laptop.png", Inches(8), Inches(1.3), width=Inches(4.8))

    tf = text(s, Inches(8), Inches(5.5), Inches(4.5), Inches(1.5), "", size=18)
    link_line(tf, "Download VS Code", "https://code.visualstudio.com", size=18)
    link_line(tf, "Download Telegram", "https://telegram.org", size=18)

    # ── SLIDE 4: WHAT PARENTS NEED TO DO ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_SOFT)
    text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.8),
         "What You Need to Do", size=36, bold=True)
    text(s, Inches(0.6), Inches(1.0), Inches(6), Inches(0.5),
         "3 Steps  ·  ~20 minutes  ·  One-time setup", size=18,
         color=ACCENT_BLUE, bold=True)
    accent_bar(s, Inches(0.6), Inches(1.5), Inches(2.5), ACCENT_BLUE)

    bullets(s, Inches(0.6), Inches(2.0), Inches(7), Inches(4.5), [
        "Step 1:  Create a Hetzner Cloud account",
        "              Cloud server for the agent — ~$4/month",
        "",
        "Step 2:  Create an OpenRouter account",
        "              AI model access — ~$1-2/month",
        "",
        "Step 3:  Install Telegram on your student's phone",
        "              Free messaging app where they talk to their agent",
        "",
        "Then share the tokens/keys with the instructor",
        "The instructor handles all the technical setup",
    ], size=21, color=TEXT_DARK, spacing=Pt(6))
    img(s, "parent_deck_3steps.png", Inches(8.2), Inches(2.0), width=Inches(4.5))

    # ── SLIDE 4: STEP 1 — HETZNER ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_WHITE)
    text(s, Inches(0.6), Inches(0.2), Inches(2), Inches(0.5),
         "STEP 1", size=16, bold=True, color=ACCENT_ORANGE)
    text(s, Inches(0.6), Inches(0.6), Inches(8), Inches(0.8),
         "Create a Hetzner Cloud Account", size=34, bold=True)
    accent_bar(s, Inches(0.6), Inches(1.3), Inches(2.5), ACCENT_ORANGE)

    bullets(s, Inches(0.6), Inches(1.7), Inches(6.8), Inches(4), [
        "1. Go to console.hetzner.cloud and sign up",
        "2. A credit card is required (billed monthly)",
        "3. Once signed in, go to:",
        "      Security  →  API Tokens  →  Generate API Token",
        "4. Select Read & Write access",
        "5. Copy the token and send it to the instructor",
        "",
        "Cost: ~$4/month for your student's server",
        "You can delete the server anytime to stop charges",
    ], size=21)

    img(s, "parent_deck_cloud.png", Inches(8), Inches(1.5), width=Inches(4.8))

    tf = text(s, Inches(8), Inches(5.5), Inches(4.5), Inches(1.5), "", size=20)
    link_line(tf, "Sign up here", "https://console.hetzner.cloud")

    # ── SLIDE 5: STEP 2 — OPENROUTER ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_WHITE)
    text(s, Inches(0.6), Inches(0.2), Inches(2), Inches(0.5),
         "STEP 2", size=16, bold=True, color=ACCENT_TEAL)
    text(s, Inches(0.6), Inches(0.6), Inches(8), Inches(0.8),
         "Create an OpenRouter Account", size=34, bold=True)
    accent_bar(s, Inches(0.6), Inches(1.3), Inches(2.5), ACCENT_TEAL)

    bullets(s, Inches(0.6), Inches(1.7), Inches(6.8), Inches(4.5), [
        "1. Go to openrouter.ai and sign up",
        "      (Google or GitHub login works)",
        "2. You get $1 in free credit immediately",
        "3. Optionally add $5 credit (lasts weeks)",
        "4. Go to Keys page  →  Create Key",
        "5. Set a credit limit on the key (e.g. $5)",
        "      This prevents overspending",
        "6. Copy the key and send to the instructor",
        "      (starts with sk-or-v1-...)",
    ], size=21)

    img(s, "parent_deck_ai.png", Inches(8), Inches(1.5), width=Inches(4.8))

    tf = text(s, Inches(0.6), Inches(6.0), Inches(11), Inches(1), "", size=18)
    link_line(tf, "Sign up", "https://openrouter.ai", size=18)
    link_line(tf, "Create key", "https://openrouter.ai/settings/keys", size=18)
    link_line(tf, "Add credit", "https://openrouter.ai/settings/credits", size=18)

    # ── SLIDE 6: STEP 3 — TELEGRAM ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_WHITE)
    text(s, Inches(0.6), Inches(0.2), Inches(2), Inches(0.5),
         "STEP 3", size=16, bold=True, color=ACCENT_BLUE)
    text(s, Inches(0.6), Inches(0.6), Inches(8), Inches(0.8),
         "Install Telegram", size=34, bold=True)
    accent_bar(s, Inches(0.6), Inches(1.3), Inches(2.5), ACCENT_BLUE)

    bullets(s, Inches(0.6), Inches(1.7), Inches(6.8), Inches(4.5), [
        "1. Download Telegram on your student's phone",
        "      Available on iOS and Android — free",
        "2. Create an account if they don't have one",
        "",
        "Your student will then do two quick things:",
        "  • Create a bot via @BotFather (2 minutes)",
        "  • Get their user ID via @userinfobot (30 seconds)",
        "",
        "The instructor will walk them through this in class",
    ], size=22)

    img(s, "parent_deck_telegram.png", Inches(8), Inches(1.3), width=Inches(4.8))

    tf = text(s, Inches(8), Inches(5.5), Inches(4.5), Inches(1.5), "", size=20)
    link_line(tf, "Download", "https://telegram.org")

    # ── SLIDE 7: COST BREAKDOWN ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_SOFT)
    text(s, Inches(0.6), Inches(0.3), Inches(8), Inches(0.8),
         "Cost Breakdown", size=36, bold=True)
    accent_bar(s, Inches(0.6), Inches(1.1), Inches(2.5), ACCENT_GREEN)

    headers = ["Item", "Monthly Cost", "Notes"]
    rows = [
        ["Cloud server (Hetzner)", "~$4/month", "Billed hourly, cancel anytime"],
        ["AI model (OpenRouter)", "~$1-2/month", "$1 free credit to start"],
        ["Telegram", "Free", ""],
        ["TOTAL", "~$5-6/month", "No contracts, cancel anytime"],
    ]

    y = Inches(1.5)
    col_lefts = [Inches(0.6), Inches(5.0), Inches(8.2)]
    col_widths = [Inches(4.2), Inches(3.0), Inches(4.5)]

    for j, val in enumerate(headers):
        text(s, col_lefts[j], y, col_widths[j], Inches(0.5),
             val, size=20, bold=True, color=ACCENT_BLUE)
    y += Inches(0.55)

    for row_data in rows:
        is_total = row_data[0] == "TOTAL"
        c = ACCENT_GREEN if is_total else TEXT_DARK
        b = is_total
        for j, val in enumerate(row_data):
            text(s, col_lefts[j], y, col_widths[j], Inches(0.5),
                 val, size=20, bold=b, color=c)
        y += Inches(0.55)

    img(s, "parent_deck_cost.png", Inches(8), Inches(3.2), width=Inches(4.5))

    text(s, Inches(0.6), Inches(5.5), Inches(7), Inches(0.8),
         "Tip: Set a credit limit on the OpenRouter key\n"
         "so your student can't accidentally overspend!",
         size=18, color=ACCENT_ORANGE, bold=True)

    # ── SLIDE 8: CHECKLIST ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_WHITE)
    text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.8),
         "Checklist: Send These to the Instructor", size=34, bold=True)
    accent_bar(s, Inches(0.6), Inches(1.1), Inches(2.5), ACCENT_GREEN)

    bullets(s, Inches(0.6), Inches(1.6), Inches(7.5), Inches(4.5), [
        "☐  Hetzner Cloud API token",
        "       (from console.hetzner.cloud → Security → API Tokens)",
        "",
        "☐  OpenRouter API key",
        "       (from openrouter.ai/settings/keys — starts with sk-or-v1-...)",
        "",
        "☐  Student's Telegram bot token",
        "       (student creates via @BotFather)",
        "",
        "☐  Student's Telegram user ID",
        "       (student gets from @userinfobot)",
        "",
        "☐  VS Code installed on student's laptop",
        "       (from code.visualstudio.com — free)",
    ], size=21, spacing=Pt(5))

    img(s, "parent_deck_checklist.png", Inches(8.5), Inches(1.3), width=Inches(4))

    text(s, Inches(0.6), Inches(6.2), Inches(10), Inches(0.6),
         "The instructor will handle all the technical setup from here.",
         size=20, color=ACCENT_GREEN, bold=True)

    # ── SLIDE 9: ALL LINKS ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_SOFT)
    text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.8),
         "All Links in One Place", size=36, bold=True)
    accent_bar(s, Inches(0.6), Inches(1.1), Inches(2.5), ACCENT_BLUE)

    tf = text(s, Inches(0.6), Inches(1.6), Inches(11), Inches(5), "", size=20)
    link_line(tf, "Hetzner Cloud (server hosting)", "https://console.hetzner.cloud", size=22)
    link_line(tf, "OpenRouter (AI model access)", "https://openrouter.ai", size=22)
    link_line(tf, "OpenRouter — Create API Key", "https://openrouter.ai/settings/keys", size=22)
    link_line(tf, "OpenRouter — Add Credits", "https://openrouter.ai/settings/credits", size=22)
    link_line(tf, "VS Code (code editor)", "https://code.visualstudio.com", size=22)
    link_line(tf, "Telegram (download app)", "https://telegram.org", size=22)

    # ── SLIDE 10: CLOSING ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_WHITE)
    text(s, Inches(0.8), Inches(2.2), Inches(11), Inches(1.5),
         "Questions?", size=48, bold=True)
    accent_bar(s, Inches(0.8), Inches(3.5), Inches(3), ACCENT_BLUE)
    text(s, Inches(0.8), Inches(4.0), Inches(11), Inches(1),
         "Reach out to the instructor if you need help with any step.",
         size=22, color=TEXT_MED)
    text(s, Inches(0.8), Inches(5.0), Inches(11), Inches(1),
         "Total setup time: ~20 minutes   |   Total cost: ~$5-6/month",
         size=20, color=ACCENT_BLUE, bold=True)

    prs.save(OUTPUT)
    print(f"Deck saved to: {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_deck()
