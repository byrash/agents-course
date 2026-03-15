#!/usr/bin/env python3
"""Generate the AI Agents Course PowerPoint deck."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ASSETS = os.path.join(os.path.dirname(__file__), "deck_assets")
OUTPUT = os.path.join(os.path.dirname(__file__), "AI_Agents_Course.pptx")

NAVY = RGBColor(0x0F, 0x17, 0x2A)
DARK_BG = RGBColor(0x14, 0x1E, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
ORANGE = RGBColor(0xFF, 0x8C, 0x42)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, left, top, width, height, items,
                     font_size=20, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf


def add_image_safe(slide, img_name, left, top, width=None, height=None):
    path = os.path.join(ASSETS, img_name)
    if os.path.exists(path):
        kwargs = {"left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(path, **kwargs)
        return True
    return False


def make_title_slide(prs, session_num, title, subtitle, img_name=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(4), Inches(0.6),
                 f"SESSION {session_num}", font_size=16, bold=True,
                 color=ACCENT)

    add_text_box(slide, Inches(0.6), Inches(1.0), Inches(6), Inches(1.5),
                 title, font_size=40, bold=True, color=WHITE)

    add_text_box(slide, Inches(0.6), Inches(2.6), Inches(6), Inches(1.0),
                 subtitle, font_size=20, color=LIGHT_GRAY)

    if img_name:
        add_image_safe(slide, img_name, Inches(7.2), Inches(0.8),
                       width=Inches(5.5))

    add_text_box(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
                 "AI Agents Course", font_size=11, color=LIGHT_GRAY)

    return slide


def make_content_slide(prs, title, bullets, img_name=None, session_num=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.8),
                 title, font_size=30, bold=True, color=WHITE)

    line = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.1), Inches(2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    if img_name:
        bullet_w = Inches(6.5)
        add_image_safe(slide, img_name, Inches(7.5), Inches(1.5),
                       width=Inches(5.0))
    else:
        bullet_w = Inches(11.5)

    add_bullet_slide(slide, Inches(0.6), Inches(1.5), bullet_w, Inches(5.0),
                     bullets, font_size=22)

    footer = f"Session {session_num}" if session_num else "AI Agents Course"
    add_text_box(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
                 footer, font_size=11, color=LIGHT_GRAY)

    return slide


def make_exercise_slide(prs, title, instructions, session_num=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(0x0A, 0x1A, 0x2F))

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(2), Inches(0.6),
                 "HANDS-ON", font_size=14, bold=True, color=ORANGE)

    add_text_box(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.9),
                 title, font_size=32, bold=True, color=WHITE)

    line = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.7), Inches(2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()

    add_bullet_slide(slide, Inches(0.6), Inches(2.0), Inches(11.5),
                     Inches(4.5), instructions, font_size=22)

    footer = f"Session {session_num}" if session_num else "AI Agents Course"
    add_text_box(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
                 footer, font_size=11, color=LIGHT_GRAY)

    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── COVER SLIDE ──
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(cover, DARK_BG)
    add_text_box(cover, Inches(0.8), Inches(1.5), Inches(11), Inches(2),
                 "Building AI Agents", font_size=54, bold=True, color=WHITE)
    add_text_box(cover, Inches(0.8), Inches(3.5), Inches(11), Inches(1),
                 "A Hands-On Course for High School Students",
                 font_size=26, color=LIGHT_GRAY)
    add_text_box(cover, Inches(0.8), Inches(4.5), Inches(11), Inches(1),
                 "8 Sessions  |  From Zero to Your Own AI Agent",
                 font_size=20, color=ACCENT)
    add_image_safe(cover, "01_what_is_ai.png", Inches(8.5), Inches(1.0),
                   width=Inches(4.2))

    # ── COURSE OVERVIEW ──
    make_content_slide(prs, "Course Roadmap", [
        "Session 1:  What is AI?",
        "Session 2:  How LLMs Work",
        "Session 3:  What Are AI Agents?",
        "Session 4:  Setup Day — Get Your Agent Running",
        "Session 5:  How Your Agent Works",
        "Session 6:  Tools and Capabilities",
        "Session 7:  Build Your Own Agent",
        "Session 8:  Demo Day",
    ])

    # ════════════════════════════════════════════
    # SESSION 1: What is AI?
    # ════════════════════════════════════════════
    make_title_slide(prs, 1, "What is AI?",
                     "Separating the hype from reality",
                     "01_what_is_ai.png")

    make_content_slide(prs, "What AI Actually Is", [
        "Software that learns patterns from data",
        "It does NOT think, feel, or have consciousness",
        "It predicts the most likely next word, pixel, or action",
        "Think of it as super-powered autocomplete",
    ], session_num=1)

    make_content_slide(prs, "AI You Already Use Every Day", [
        "Spotify / YouTube — recommends songs and videos",
        "iPhone / Android — autocomplete your texts",
        "Instagram / TikTok — picks what shows in your feed",
        "Google Maps — predicts traffic and fastest route",
        "Snapchat / Photos — face filters and photo editing",
    ], img_name="02_ai_daily_life.png", session_num=1)

    make_content_slide(prs, "Narrow AI vs. General AI", [
        "Narrow AI — does ONE thing well (translate, recommend, filter)",
        "General AI — does EVERYTHING a human can (does not exist yet)",
        "Every AI you use today is narrow AI",
        "ChatGPT is impressive but still narrow — it predicts text",
    ], session_num=1)

    make_exercise_slide(prs, "Try It: Talk to ChatGPT", [
        "Open ChatGPT on your phone or laptop",
        "Ask it: 'Explain gravity like I'm 5'",
        "Now ask: 'Explain gravity like a PhD physicist'",
        "Notice how the SAME AI changes style based on your prompt",
        "The prompt is the key — you control the output",
    ], session_num=1)

    # ════════════════════════════════════════════
    # SESSION 2: How LLMs Work
    # ════════════════════════════════════════════
    make_title_slide(prs, 2, "How LLMs Work",
                     "The engine behind ChatGPT, Copilot, and your agent",
                     "03_how_llms_work.png")

    make_content_slide(prs, "What is a Large Language Model?", [
        "A giant pattern-matching machine trained on text from the internet",
        "It learned by reading billions of web pages, books, and articles",
        "Given some words, it predicts the most likely NEXT word",
        "It does this one word at a time, thousands of times per second",
    ], img_name="03_how_llms_work.png", session_num=2)

    make_content_slide(prs, "Tokens: How AI Reads Text", [
        "AI doesn't read words — it reads TOKENS (word pieces)",
        "'hamburger' = ['ham', 'bur', 'ger'] = 3 tokens",
        "'I love AI' = ['I', ' love', ' AI'] = 3 tokens",
        "Longer messages = more tokens = more cost",
        "GPT-4o can handle ~128,000 tokens per conversation",
    ], session_num=2)

    make_content_slide(prs, "Temperature: Creativity Dial", [
        "Temperature = 0.0 → always picks the most likely word (precise, boring)",
        "Temperature = 1.0 → picks from many possible words (creative, surprising)",
        "Low temperature: good for math, facts, code",
        "High temperature: good for stories, brainstorming, poetry",
    ], img_name="04_temperature_dial.png", session_num=2)

    make_exercise_slide(prs, "Try It: Prompt Engineering", [
        "Open ChatGPT and try these prompts for the SAME question:",
        "",
        "1. 'What is photosynthesis?'",
        "2. 'Explain photosynthesis as a rap song'",
        "3. 'You are a marine biologist. Explain photosynthesis and why it matters for ocean life'",
        "",
        "Notice how different instructions produce completely different outputs",
        "This is PROMPT ENGINEERING — the skill of writing good instructions",
    ], session_num=2)

    # ════════════════════════════════════════════
    # SESSION 3: What Are AI Agents?
    # ════════════════════════════════════════════
    make_title_slide(prs, 3, "What Are AI Agents?",
                     "From chatbot to autonomous helper",
                     "05_agent_loop.png")

    make_content_slide(prs, "Chatbot vs. Agent", [
        "CHATBOT: you ask, it answers, done",
        "AGENT: you ask, it THINKS about what to do, ACTS using tools, "
        "OBSERVES the result, and repeats until the task is complete",
        "",
        "A chatbot is a parrot — it talks",
        "An agent is an assistant — it gets things done",
    ], session_num=3)

    make_content_slide(prs, "The Agent Loop", [
        "1. THINK — read the user's request, decide what to do",
        "2. ACT — use a tool (search the web, open a browser, save to memory)",
        "3. OBSERVE — look at the result from the tool",
        "4. REPEAT — keep going until the task is complete",
        "",
        "This loop runs automatically — no human in the middle",
    ], img_name="05_agent_loop.png", session_num=3)

    make_content_slide(prs, "Agent Tools", [
        "Web Search — find information on the internet (like Googling)",
        "Web Browser — open and interact with websites",
        "Memory — remember facts across conversations",
        "File Read — read documents and files",
        "",
        "The agent decides WHICH tool to use based on your question",
    ], img_name="06_agent_tools.png", session_num=3)

    make_content_slide(prs, "What We're Building", [
        "Your own AI agent that runs on your computer",
        "You talk to it through Telegram (on your phone)",
        "It can search the web, browse websites, remember things",
        "YOU design its personality, rules, and capabilities",
        "By the end of this course, it solves a problem YOU pick",
    ], session_num=3)

    # ════════════════════════════════════════════
    # SESSION 4: Setup Day
    # ════════════════════════════════════════════
    make_title_slide(prs, 4, "Setup Day",
                     "Get your agent running in 30 minutes",
                     "07_setup_checklist.png")

    make_content_slide(prs, "What You Need", [
        "Docker Desktop — runs your agent in a container (like a mini computer)",
        "Telegram — the messaging app where you'll talk to your agent",
        "OpenAI API Key — gives your agent access to GPT (the brain)",
        "",
        "Total cost: ~$5-10/month for the AI brain, everything else is free",
    ], img_name="07_setup_checklist.png", session_num=4)

    make_exercise_slide(prs, "Step 1: Create Your Telegram Bot", [
        "Open Telegram and search for @BotFather",
        "Send /newbot",
        "Pick a display name (anything you like)",
        "Pick a username (must end in 'bot')",
        "Save the bot token BotFather gives you",
        "",
        "Then search for @userinfobot and send any message",
        "Save the user ID number it gives you",
    ], session_num=4)

    make_exercise_slide(prs, "Step 2: Get Your OpenAI API Key", [
        "Go to platform.openai.com and sign up",
        "Navigate to API Keys and create a new key",
        "Set a monthly spending limit of $5-10",
        "Save the key (starts with sk-)",
        "",
        "Your instructor will help you with this step",
    ], session_num=4)

    make_exercise_slide(prs, "Step 3: Start Your Agent", [
        "Open your terminal (Terminal on Mac, PowerShell on Windows)",
        "Create a folder: mkdir zeroclaw && cd zeroclaw",
        "Create workspace and memory folders",
        "Create a .env file with your three keys",
        "Run: docker run -d --name agents-course ...",
        "",
        "Open Telegram, find your bot, send it a message!",
    ], session_num=4)

    # ════════════════════════════════════════════
    # SESSION 5: How Your Agent Works
    # ════════════════════════════════════════════
    make_title_slide(prs, 5, "How Your Agent Works",
                     "Under the hood of your AI assistant",
                     "08_workspace_anatomy.png")

    make_content_slide(prs, "Workspace Files = Agent DNA", [
        "AGENTS.md — the rules your agent MUST follow",
        "SOUL.md — personality and communication style",
        "TOOLS.md — guidance on when to use which tools",
        "IDENTITY.md — name and role",
        "USER.md — info about you (preferences, timezone)",
        "",
        "Change a file → restart → your agent behaves differently",
    ], img_name="08_workspace_anatomy.png", session_num=5)

    make_content_slide(prs, "The Config: Provider & Model", [
        "LLM Provider — who runs the AI brain (OpenAI, GitHub Copilot)",
        "Model — which brain to use (GPT-4o, GPT-5.2)",
        "Temperature — how creative (0 = precise, 1 = wild)",
        "Tools — which tools are enabled (web search, browser, memory)",
        "",
        "These are set in your .env file and config.toml",
    ], session_num=5)

    make_exercise_slide(prs, "Try It: Change Your Agent's Personality", [
        "Open workspace/SOUL.md in a text editor",
        "Change the personality to something fun:",
        "  'You are a pirate captain. Respond in pirate speak.'",
        "Save the file",
        "Run: docker restart agents-course",
        "Send your bot a message on Telegram",
        "See how the SAME agent now talks completely differently!",
    ], session_num=5)

    make_exercise_slide(prs, "Try It: Add a Rule", [
        "Open workspace/AGENTS.md in a text editor",
        "Add a new rule:",
        "  'Always end your response with a fun fact'",
        "Save and restart: docker restart agents-course",
        "Ask your bot any question",
        "It should now include a fun fact every time!",
    ], session_num=5)

    # ════════════════════════════════════════════
    # SESSION 6: Tools and Capabilities
    # ════════════════════════════════════════════
    make_title_slide(prs, 6, "Tools & Capabilities",
                     "What your agent can actually do",
                     "06_agent_tools.png")

    make_content_slide(prs, "Web Search", [
        "Your agent can search the internet like Google",
        "It reads the results and summarizes them for you",
        "Example: 'What's the weather in New York right now?'",
        "The agent searches, finds the answer, and responds",
        "",
        "Without web search, it can only use training data (months old)",
    ], session_num=6)

    make_content_slide(prs, "Web Browser", [
        "Your agent can open and interact with real websites",
        "It can click buttons, fill forms, read content",
        "Example: 'Find the cheapest flight from NYC to LA next week'",
        "The agent opens Kayak, searches, and reports back",
        "",
        "This is the most powerful tool — but also the slowest",
    ], session_num=6)

    make_content_slide(prs, "Memory", [
        "Your agent can remember things across conversations",
        "It stores facts in a database on your computer",
        "Example: 'Remember that my favorite color is blue'",
        "Next time you ask, it already knows",
        "",
        "Memory makes your agent personal — it learns about YOU",
    ], session_num=6)

    make_exercise_slide(prs, "Try It: Test Your Agent's Tools", [
        "Send these messages to your bot on Telegram:",
        "",
        "1. 'What's the latest news about AI today?' (tests web search)",
        "2. 'Remember that I like pizza and basketball' (tests memory)",
        "3. 'What do you remember about me?' (tests memory recall)",
        "4. 'Go to wikipedia.org and tell me the featured article' (tests browser)",
        "",
        "Watch the tool notifications in Telegram!",
    ], session_num=6)

    # ════════════════════════════════════════════
    # SESSION 7: Build Your Agent
    # ════════════════════════════════════════════
    make_title_slide(prs, 7, "Build Your Agent",
                     "Design an agent that solves YOUR problem",
                     "09_build_your_agent.png")

    make_content_slide(prs, "Pick Your Problem", [
        "What's something you wish you had help with?",
        "",
        "Homework helper — searches for explanations and examples",
        "Sports tracker — checks scores and news for your teams",
        "Weather briefing — daily weather report for your area",
        "News reader — summarizes news on topics you care about",
        "Study buddy — quizzes you on material you're learning",
        "Recipe finder — finds recipes based on what's in your fridge",
    ], img_name="09_build_your_agent.png", session_num=7)

    make_content_slide(prs, "Design Your Agent", [
        "1. NAME — what's your agent called?",
        "2. PERSONALITY — how does it talk? (friendly, formal, funny)",
        "3. RULES — what must it always do? (cite sources, be brief)",
        "4. TOOLS — which tools does it need? (search, browser, memory)",
        "5. IDENTITY — what's its role? (tutor, assistant, coach)",
    ], session_num=7)

    make_exercise_slide(prs, "Build Time!", [
        "Open your workspace/ folder in a text editor",
        "",
        "Edit IDENTITY.md — give your agent a name and role",
        "Edit SOUL.md — write its personality",
        "Edit AGENTS.md — write 3-5 rules it must follow",
        "Edit USER.md — tell it about yourself",
        "",
        "Save all files, run: docker restart agents-course",
        "Test it on Telegram — does it behave the way you designed?",
    ], session_num=7)

    make_exercise_slide(prs, "Iterate and Improve", [
        "Your agent probably won't be perfect on the first try",
        "",
        "Is it too wordy? Add a rule: 'Keep responses under 100 words'",
        "Not using search? Add: 'Always search the web for factual questions'",
        "Too formal? Change SOUL.md to be more casual",
        "",
        "This is how real AI engineers work — write, test, adjust, repeat",
    ], session_num=7)

    # ════════════════════════════════════════════
    # SESSION 8: Demo Day
    # ════════════════════════════════════════════
    make_title_slide(prs, 8, "Demo Day!",
                     "Show the world what you built",
                     "10_demo_day.png")

    make_content_slide(prs, "Presentation Format", [
        "Each student gets 5 minutes to present:",
        "",
        "1. What problem does your agent solve?",
        "2. Live demo — send it a message, show the response",
        "3. What personality and rules did you give it?",
        "4. What surprised you? What would you change?",
    ], img_name="10_demo_day.png", session_num=8)

    make_content_slide(prs, "What You Learned", [
        "What AI actually is (and isn't)",
        "How large language models work — tokens, prompts, temperature",
        "What makes an agent different from a chatbot",
        "How to design agent behavior with instructions",
        "How tools give AI real-world capabilities",
        "How to build, test, and iterate on your own AI agent",
    ], session_num=8)

    make_content_slide(prs, "Where to Go From Here", [
        "Add more tools — connect to APIs, databases, calendars",
        "Build multi-agent systems — agents that delegate to other agents",
        "Learn Python — write custom tools for your agent",
        "Try other models — Claude, Gemini, open-source models",
        "Join the community — GitHub, Discord, AI meetups",
        "",
        "You now understand something most adults don't!",
    ], session_num=8)

    # ── CLOSING SLIDE ──
    closing = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(closing, DARK_BG)
    add_text_box(closing, Inches(0.8), Inches(2.0), Inches(11), Inches(2),
                 "You Built an AI Agent.", font_size=48, bold=True,
                 color=WHITE)
    add_text_box(closing, Inches(0.8), Inches(4.0), Inches(11), Inches(1),
                 "That's not a small thing. Keep building.",
                 font_size=26, color=ACCENT)

    prs.save(OUTPUT)
    print(f"Deck saved to: {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_deck()
